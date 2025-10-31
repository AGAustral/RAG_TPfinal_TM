"""
Script simple para convertir PDF, DOCX, PPTX, XLSX a JSONL
Uso: Coloca tus archivos en una carpeta y ejecuta este script
"""

import os
import json
from pathlib import Path
from datetime import datetime
import warnings
warnings.filterwarnings("ignore")

# Instalar si no tienes: pip install PyPDF2 python-docx python-pptx openpyxl
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl


def extraer_texto_pdf(filepath):
    """Extrae texto de PDF"""
    texto = ""
    try:
        with open(filepath, 'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            for pagina in pdf.pages:
                texto += pagina.extract_text() + "\n"
    except Exception as e:
        print(f"Error en PDF {filepath}: {e}")
    return texto.strip()


def extraer_texto_docx(filepath):
    """Extrae texto de DOCX"""
    texto = ""
    try:
        doc = Document(filepath)
        texto = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"Error en DOCX {filepath}: {e}")
    return texto.strip()


def extraer_texto_pptx(filepath):
    """Extrae texto de PPTX/PPTM"""
    texto = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
    except Exception as e:
        print(f"Error en PPTX {filepath}: {e}")
    return texto.strip()


def extraer_texto_xlsx(filepath):
    """Extrae texto de XLSX"""
    texto = ""
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        for hoja in wb.sheetnames:
            sheet = wb[hoja]
            texto += f"\n[Hoja: {hoja}]\n"
            for fila in sheet.iter_rows(values_only=True):
                fila_texto = " | ".join([str(c) for c in fila if c is not None])
                if fila_texto.strip():
                    texto += fila_texto + "\n"
    except Exception as e:
        print(f"Error en XLSX {filepath}: {e}")
    return texto.strip()


def dividir_en_chunks(texto, nombre_archivo, palabras_por_chunk=350, overlap=150):
    """Divide texto en chunks: primero por separador temático, luego por longitud."""
    
    chunks_finales = []
    
    # 1. DIVISIÓN TEMÁTICA: Usar el separador lógico si existe
    if '$$$' in texto:
        bloques_tematicos = [b.strip() for b in texto.split('$$$') if b.strip()]
    else:
        # Si no hay separador, usar el texto completo como un solo bloque
        bloques_tematicos = [texto]

    # 2. PROCESAR CADA BLOQUE
    for bloque_texto in bloques_tematicos:
        palabras = bloque_texto.split()
        
        # 3. DIVISIÓN POR LONGITUD (Fallback)
        # Esto solo se ejecutará si un bloque temático es mayor que 500 palabras
        for i in range(0, len(palabras), palabras_por_chunk - overlap):
            chunk_palabras = palabras[i:i + palabras_por_chunk]
            chunk_texto = " ".join(chunk_palabras)
            
            if chunk_texto.strip():
                chunks_finales.append({
                    "id": f"{Path(nombre_archivo).stem}_chunk_{len(chunks_finales)}",
                    "text": chunk_texto.strip(),
                    "metadata": {
                        "source": nombre_archivo,
                        "chunk_index": len(chunks_finales),
                        "date_processed": datetime.now().isoformat()
                    }
                })
    
    # 4. Actualizar total de chunks
    total = len(chunks_finales)
    for chunk in chunks_finales:
        chunk["metadata"]["total_chunks"] = total
    
    return chunks_finales


def procesar_archivo(filepath):
    """Procesa un archivo según su extensión"""
    ext = Path(filepath).suffix.lower()
    nombre = Path(filepath).name
    
    print(f"Procesando: {nombre}")
    
    # Extraer texto según formato
    if ext == '.pdf':
        texto = extraer_texto_pdf(filepath)
    elif ext == '.docx':
        texto = extraer_texto_docx(filepath)
    elif ext in ['.pptx', '.pptm']:
        texto = extraer_texto_pptx(filepath)
    elif ext == '.xlsx':
        texto = extraer_texto_xlsx(filepath)
    else:
        print(f"⚠️  Formato no soportado: {ext}")
        return []
    
    if not texto:
        print(f"⚠️  Sin texto extraído de {nombre}")
        return []
    
    # Dividir en chunks
    chunks = dividir_en_chunks(texto, nombre)
    print(f"✓ {len(chunks)} chunks generados")
    
    return chunks


def procesar_carpeta(carpeta_entrada, archivo_salida, palabras_por_chunk=350):
    """
    Procesa todos los archivos de una carpeta y guarda en JSONL
    
    Args:
        carpeta_entrada: Carpeta con tus archivos (PDF, DOCX, etc.)
        archivo_salida: Archivo JSONL de salida
        palabras_por_chunk: Tamaño de cada chunk en palabras
    """
    
    # Buscar archivos
    extensiones = ['.pdf', '.docx', '.pptx', '.pptm', '.xlsx']
    archivos = [f for f in Path(carpeta_entrada).glob('*') 
                if f.suffix.lower() in extensiones]
    
    if not archivos:
        print(f"⚠️  No se encontraron archivos en {carpeta_entrada}")
        return
    
    print(f"\n{'='*60}")
    print(f"Procesando {len(archivos)} archivos")
    print(f"{'='*60}\n")
    
    todos_chunks = []
    
    # Procesar cada archivo
    for archivo in archivos:
        chunks = procesar_archivo(str(archivo))
        todos_chunks.extend(chunks)
    
    # Guardar JSONL
    Path(archivo_salida).parent.mkdir(parents=True, exist_ok=True)
    
    with open(archivo_salida, 'w', encoding='utf-8') as f:
        for chunk in todos_chunks:
            f.write(json.dumps(chunk, ensure_ascii=False) + '\n')
    
    print(f"\n{'='*60}")
    print(f"✓ Completado!")
    print(f"✓ Total chunks: {len(todos_chunks)}")
    print(f"✓ Guardado en: {archivo_salida}")
    print(f"{'='*60}\n")
    
    # Resumen por archivo
    print("Resumen por archivo:")
    resumen = {}
    for chunk in todos_chunks:
        source = chunk['metadata']['source']
        resumen[source] = resumen.get(source, 0) + 1
    
    for archivo, cantidad in resumen.items():
        print(f"  • {archivo}: {cantidad} chunks")


# ==================== CONFIGURACIÓN Y USO ====================

if __name__ == "__main__":
    

    # CONFIGURA ESTAS RUTAS
    CARPETA_ENTRADA = r"C:\Users\agarmendia\Desktop\PNCIL_TinyLlama"
    ARCHIVO_SALIDA = r"C:\Users\agarmendia\Desktop\PNCIL_TinyLlama\pncil_documents.jsonl"
    PALABRAS_POR_CHUNK = 500                  # Tamaño de chunks
    
    # Procesar
    procesar_carpeta(CARPETA_ENTRADA, ARCHIVO_SALIDA, PALABRAS_POR_CHUNK)
    
    print("\n✓ Ahora puedes usar este JSONL en tu código RAG original")